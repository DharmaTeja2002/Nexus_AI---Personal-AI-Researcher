from pathlib import Path
from typing import List
from docx import Document # python-docx
from pptx import Presentation # python-pptx
from src.ingestion.schemas import DocumentElement

class OfficeParser:
    """
    Worker responsible for extracting content from MS Word (.docx) 
    and MS PowerPoint (.pptx) files.
    """
    def parse(self, path: Path) -> List[DocumentElement]:
        # Determine if it's Word or PowerPoint based on extension
        if path.suffix.lower() == '.docx':
            return self._parse_word(path)
        elif path.suffix.lower() == '.pptx':
            return self._parse_pptx(path)
        return []

    def _parse_word(self, path: Path) -> List[DocumentElement]:
        try:
            doc = Document(path)
            elements = []
            # Extract every paragraph
            for i, para in enumerate(doc.paragraphs):
                if para.text.strip():
                    elements.append(
                        DocumentElement(
                            element_type="text",
                            content=para.text.strip(),
                            source_file=path.name,
                            metadata={"paragraph_index": i, "format": "docx"}
                        )
                    )
            return elements
        except Exception as e:
            print(f"❌ Error parsing Word file {path.name}: {e}")
            return []

    def _parse_pptx(self, path: Path) -> List[DocumentElement]:
        try:
            prs = Presentation(path)
            elements = []
            # Iterate through slides and then shapes in each slide
            for s_idx, slide in enumerate(prs.slides):
                for sh_idx, shape in enumerate(slide.shapes):
                    if hasattr(shape, "text") and shape.text.strip():
                        elements.append(
                            DocumentElement(
                                element_type="text",
                                content=shape.text.strip(),
                                source_file=path.name,
                                metadata={"slide_index": s_idx, "shape_index": sh_idx, "format": "pptx"}
                            )
                        )
            return elements
        except Exception as e:
            print(f"❌ Error parsing PPTX file {path.name}: {e}")
            return []